#!/usr/bin/env python3
"""スマートホーム手順書のデザイン・トーンを踏襲した入居者向けスマートロック手順書を生成する。"""
from __future__ import annotations

import copy
import shutil
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(
    "/Users/matsunomasaharu2/Library/CloudStorage/OneDrive-個人用/215_神・大家さん倶楽部/"
    "20_【空室対策】【修繕】【売却】/21_【空室対策】募集,ステージング,物件管理"
)
TEMPLATE = BASE_DIR / "スマートホーム手順書【入居者向け】.pptx"
OUT = BASE_DIR / "スマートロック手順書【入居者向け】.pptx"
ASSETS_DIR = BASE_DIR / "_sesami_assets"

# CANDY HOUSE 公式ストア（Shopify）の製品画像
ASSET_URLS: dict[str, str] = {
    "sesame5_main.png": (
        "https://cdn.shopify.com/s/files/1/0016/1870/6495/files/52.png?v=1747630869"
    ),
    "sesame5_lifestyle.png": (
        "https://cdn.shopify.com/s/files/1/0016/1870/6495/files/"
        "30_a53b2692-9c66-4000-ad18-711cea5be676.png?v=1756798270"
    ),
    "sesame5_app.png": (
        "https://cdn.shopify.com/s/files/1/0016/1870/6495/files/332.png?v=1747630869"
    ),
    "sesame_touch2_main.png": (
        "https://cdn.shopify.com/s/files/1/0016/1870/6495/files/SSMTouch_L.png?v=1748840832"
    ),
    "sesame_touch2_hero.png": (
        "https://cdn.shopify.com/s/files/1/0016/1870/6495/files/hp_touch.png?v=1770786842"
    ),
}

# テンプレート内スライド番号（1始まり）— 削除前に参照
IDX_TITLE = 0
IDX_INTRO = 1
IDX_DEVICE_A = 2
IDX_DEVICE_B = 3
IDX_DIAGRAM = 4
IDX_SECTION = 5
DELETE_FROM = 6  # 旧スライド7〜を削除
IDX_TRY_ORIG = 21
IDX_LIST_ORIG = 22


def prep_templates(prs: Presentation) -> dict[str, int]:
    """削除後も使えるよう、セクション・リスト用レイアウトを末尾に退避。"""
    duplicate_slide(prs, IDX_SECTION)
    duplicate_slide(prs, IDX_TRY_ORIG)
    duplicate_slide(prs, IDX_LIST_ORIG)
    last_to_delete = len(prs.slides) - 4  # 退避3枚の直前まで
    for i in range(last_to_delete, DELETE_FROM - 1, -1):
        delete_slide(prs, i)
    # 削除後は 0-5=イントロ、6-8=複製テンプレート
    return {"section": 6, "try": 7, "list": 8}


def cleanup_templates(prs: Presentation, tpl: dict[str, int]) -> None:
    for idx in sorted(tpl.values(), reverse=True):
        delete_slide(prs, idx)


def duplicate_slide(prs: Presentation, index: int):
    template = prs.slides[index]
    copied = prs.slides.add_slide(template.slide_layout)
    for shp in template.shapes:
        newel = copy.deepcopy(shp.element)
        copied.shapes._spTree.insert_element_before(newel, "p:extLst")
    return copied


def delete_slide(prs: Presentation, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def ensure_assets() -> None:
    ASSETS_DIR.mkdir(exist_ok=True)
    for name, url in ASSET_URLS.items():
        path = ASSETS_DIR / name
        if path.exists() and path.stat().st_size > 1000:
            continue
        data = urllib.request.urlopen(url, timeout=60).read()
        path.write_bytes(data)


def picture_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def remove_pictures(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in picture_shapes(slide):
        sp_tree.remove(shape._element)


def swap_pictures(slide, image_paths: list[Path]) -> None:
    """既存写真の位置・サイズを保ったまま、上から順に差し替える。"""
    boxes = sorted(
        [(s.left, s.top, s.width, s.height) for s in picture_shapes(slide)],
        key=lambda b: b[2],
        reverse=True,
    )
    remove_pictures(slide)
    for box, path in zip(boxes[: len(image_paths)], image_paths):
        slide.shapes.add_picture(str(path), box[0], box[1], width=box[2], height=box[3])


def patch_diagram_pictures(slide) -> None:
    large = [s for s in picture_shapes(slide) if s.width >= 2000000 and s.height >= 2000000]
    boxes = sorted(
        [(s.left, s.top, s.width, s.height) for s in large],
        key=lambda b: b[0],
    )
    sp_tree = slide.shapes._spTree
    for shape in picture_shapes(slide):
        sp_tree.remove(shape._element)
    paths = [ASSETS_DIR / "sesame_touch2_main.png", ASSETS_DIR / "sesame5_main.png"]
    for box, path in zip(boxes[: len(paths)], paths):
        slide.shapes.add_picture(str(path), box[0], box[1], width=box[2], height=box[3])


def set_shape_lines(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, line in enumerate(lines):
        p = tf.paragraphs[i]
        if p.runs:
            p.runs[0].text = line
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = line
    for j in range(len(lines), len(tf.paragraphs)):
        p = tf.paragraphs[j]
        p.text = ""
        for r in p.runs:
            r.text = ""


def clear_shape_text(shape) -> None:
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        p.text = ""
        for r in p.runs:
            r.text = ""


def set_single_line(shape, text: str) -> None:
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = text


def first_text_shape(slide, skip: int = 0):
    found = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            if found == skip:
                return shape
            found += 1
    return None


def list_slide(prs: Presentation, tpl: dict[str, int], title: str, columns: list[list[str]]):
    slide = duplicate_slide(prs, tpl["list"])
    title_shape = None
    col_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if "たとえば" in t or "こんなこと" in t:
            title_shape = shape
        elif t.startswith("「アレクサ"):
            col_shapes.append(shape)
    col_shapes.sort(key=lambda s: (s.top, s.left))
    if title_shape:
        set_single_line(title_shape, title)
    for shape, lines in zip(col_shapes[:3], columns):
        set_shape_lines(shape, lines)
    return slide


def section_slide(prs: Presentation, tpl: dict[str, int], title: str):
    slide = duplicate_slide(prs, tpl["section"])
    set_single_line(slide.shapes[0], title)
    return slide


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame]


def patch_intro_slides(prs: Presentation) -> None:
    s1 = prs.slides[IDX_TITLE]
    set_single_line(s1.shapes[1], "スマートロックのつくりかた")
    set_single_line(s1.shapes[2], "※鍵の設定は、入居日に大家（オーナー）と一緒に行ってください。")

    s2 = prs.slides[IDX_INTRO]
    # 全面背景写真（sh0）とグレーパネルはテンプレートのまま残す
    intro = next(
        (s for s in s2.shapes if s.has_text_frame and "デバイス" in s.text_frame.text),
        first_text_shape(s2),
    )
    if intro:
        set_shape_lines(
            intro,
            [
                "鍵を持たずに、スマホや指紋で開けられる！",
                "そんなスマートロックは、２つのデバイスの",
                "コンビネーションでつくることができます。",
                "２つのデバイスの役割をそれぞれ理解していると",
                "仕組みがわかりやすいので、まずはそのご説明から。",
            ],
        )

    s3 = prs.slides[IDX_DEVICE_A]
    swap_pictures(
        s3,
        [ASSETS_DIR / "sesame5_lifestyle.png", ASSETS_DIR / "sesame5_app.png"],
    )
    ts3 = text_shapes(s3)
    set_single_line(ts3[0], "SESAME 5")
    set_shape_lines(
        ts3[1],
        [
            "ドアのサムターンを回して、実際に施錠・解錠する本体です。",
            "スマホのセサミアプリから、Bluetoothで操作できます。",
            "（玄関のすぐ近くで使います）",
        ],
    )

    s4 = prs.slides[IDX_DEVICE_B]
    swap_pictures(s4, [ASSETS_DIR / "sesame_touch2_hero.png"])
    ts4 = text_shapes(s4)
    set_single_line(ts4[0], "SESAMEタッチ2")
    set_shape_lines(
        ts4[1],
        [
            "指紋・ICカード・スマホで「開けて」と指示する認証パネルです。",
            "タッチするだけで、SESAME 5に解錠の指示を送れます。",
            "（SESAMEタッチ2だけでは鍵は開きません）",
        ],
    )

    s5 = prs.slides[IDX_DIAGRAM]
    patch_diagram_pictures(s5)
    title5 = next((s for s in text_shapes(s5) if "仕組み" in s.text_frame.text or "スマート" in s.text_frame.text), text_shapes(s5)[0])
    set_single_line(title5, "つまり、スマートロックの仕組みはこういうこと！")
    replacements = [
        ("声に反応", ["スマホで", "セサミアプリを", "操作。"]),
        ("対応する家電", ["Bluetoothで", "SESAME 5に", "伝える"]),
        ("パッ！", ["カチッ"]),
        ("アレクサ、\nテレビ", ["指紋や", "ICカードで", "タッチ！"]),
        ("アレクサ、\nエアコン", ["スマホを", "かざして", "タッチ！"]),
        ("アレクサ、\n電気", ["SESAME 5が", "ドアの鍵を", "回す"]),
        ("wifi", ["Bluetoothでつながっている"]),
        ("声でアレクサ", ["スマホのセサミアプリからSESAME 5を操作。", "SESAMEタッチ2は指紋・IC・スマホでSESAME 5に解錠指示を送る。"]),
    ]
    for key, lines in replacements:
        for shape in text_shapes(s5):
            if key in shape.text_frame.text.replace("\n", ""):
                set_shape_lines(shape, lines)
                break
    stale = ("アレクサ", "Nature", "Remo", "パッ！", "ピッ！", "テレビ", "エアコン", "電気つけ")
    for shape in text_shapes(s5):
        if shape is title5:
            continue
        t = shape.text_frame.text
        if any(k in t for k in stale):
            clear_shape_text(shape)

    s6 = prs.slides[IDX_SECTION]
    set_single_line(s6.shapes[0], "入居者さまの準備のしかた")


def build_content_slides(prs: Presentation, tpl: dict[str, int]) -> None:
    list_slide(
        prs,
        tpl,
        "まず、セサミアプリを入れましょう",
        [
            [
                "1. App Store または",
                "Google Play で",
                "「SESAME」を検索",
            ],
            [
                "2. セサミアプリを",
                "インストール",
            ],
            [
                "3. 必要環境",
                "Bluetooth 4.0以上",
                "iOS 8以降 / Android 4.3以降",
            ],
        ],
    )

    list_slide(
        prs,
        tpl,
        "メールアドレスは、大家に伝えたもので！",
        [
            [
                "大家にあらかじめ",
                "伝えたメールアドレスで",
                "登録・ログインしてください",
            ],
            [
                "別のアドレスで登録すると",
                "鍵が表示されません",
            ],
            [
                "Facebook / Google 連携は",
                "登録後に追加できます",
            ],
        ],
    )

    section_slide(prs, tpl, "鍵の受け取り方")

    list_slide(
        prs,
        tpl,
        "メールで鍵を受け取る場合",
        [
            [
                "1. セサミアプリを開く",
                "2. ログイン済みなら",
                "お部屋のSESAMEが",
                "一覧に表示されます",
            ],
            [
                "表示されないときは",
                "メールアドレスが",
                "一致しているか確認",
            ],
            [
                "大家がシェア後",
                "数分待ってから",
                "再度アプリを開いてみてください",
            ],
        ],
    )

    list_slide(
        prs,
        tpl,
        "QRコードで鍵を受け取る場合",
        [
            [
                "1. セサミアプリを開く",
                "2. 右上の「⊕」をタップ",
                "3.「QRコードのスキャン」",
            ],
            [
                "大家から受け取った",
                "QRコードを読み取る",
            ],
            [
                "お部屋のSESAMEが",
                "追加されます",
            ],
        ],
    )

    try_slide = duplicate_slide(prs, tpl["try"])
    set_single_line(try_slide.shapes[0], "使ってみましょう！")

    list_slide(
        prs,
        tpl,
        "スマホで解錠・施錠するには",
        [
            [
                "1. 玄関のすぐ近くで",
                "セサミアプリを開く",
                "（Bluetooth圏内）",
            ],
            [
                "2. お部屋のSESAMEを",
                "タップ",
            ],
            [
                "3.「開ける」「閉める」",
                "ボタンを操作",
                "※遠隔操作はマネージャー権限＋",
                "大家のWi-Fi設置時のみ",
            ],
        ],
    )

    list_slide(
        prs,
        tpl,
        "便利機能（お好みで設定）",
        [
            [
                "手ぶら解錠",
                "近づくと自動で解錠",
            ],
            [
                "ノック解錠",
                "ドアをノックすると解錠",
            ],
            [
                "マネージャー・ゲスト",
                "どちらも設定できます",
            ],
        ],
    )

    list_slide(
        prs,
        tpl,
        "SESAMEタッチ2で解錠（指紋）",
        [
            [
                "登録済みの指を",
                "タッチパネルの",
                "ボタン部分に当てる",
            ],
            [
                "登録は入居時に",
                "大家と一緒に行います",
            ],
            [
                "追加したいときは",
                "大家に連絡してください",
            ],
        ],
    )

    list_slide(
        prs,
        tpl,
        "SESAMEタッチ2で解錠（IC・スマホ）",
        [
            [
                "ICカードを",
                "タッチパネルにかざす",
            ],
            [
                "スマホ（Felica/NFC）を",
                "タッチパネルにかざす",
            ],
            [
                "登録は入居時に",
                "大家と一緒に行います",
            ],
        ],
    )

    section_slide(prs, tpl, "困ったとき")

    list_slide(
        prs,
        tpl,
        "うまくいかないときは",
        [
            [
                "鍵が表示されない",
                "→ 大家に伝えたメールで",
                "ログインしているか確認",
            ],
            [
                "スマホで開けない",
                "→ 玄関に近づく",
                "BluetoothをON",
            ],
            [
                "タッチ2が反応しない",
                "→ 電池・指紋登録を大家に確認",
                "→ 物理鍵（バックアップ鍵）を使用",
            ],
        ],
    )

    list_slide(
        prs,
        tpl,
        "お問い合わせ",
        [
            [
                "大家（オーナー）",
                "松野",
                "090-9670-7595",
            ],
            [
                "セサミサポート",
                "（CANDY HOUSE）",
                "sesame@candyhouse.co",
            ],
            [
                "ヘルプセンター",
                "jp.candyhouse.co",
                "/apps/help-center",
            ],
        ],
    )


def main() -> None:
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")

    ensure_assets()
    shutil.copy2(TEMPLATE, OUT)
    prs = Presentation(str(OUT))

    patch_intro_slides(prs)
    tpl = prep_templates(prs)
    build_content_slides(prs, tpl)
    cleanup_templates(prs, tpl)

    prs.save(str(OUT))
    print(f"Wrote: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
